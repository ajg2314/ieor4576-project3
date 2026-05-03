from io import BytesIO
import pymupdf


def extract_text_from_pdf(file_bytes: bytes) -> str:
    doc = pymupdf.open(stream=file_bytes, filetype="pdf")
    return "\n\n".join(page.get_text() for page in doc)


def extract_text_from_docx(file_bytes: bytes) -> str:
    import docx
    doc = docx.Document(BytesIO(file_bytes))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def clean_text(raw: str) -> str:
    lines = [line.strip() for line in raw.splitlines()]
    result, prev_blank = [], False
    for line in lines:
        if not line:
            if not prev_blank:
                result.append("")
            prev_blank = True
        else:
            result.append(line)
            prev_blank = False
    return "\n".join(result).strip()
